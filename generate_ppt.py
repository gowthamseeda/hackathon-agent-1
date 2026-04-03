from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # --- Slide 1: Title Slide ---
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Customer Feedback Intelligence AI"
    subtitle.text = "Orchestrating Insights with Gemini 2.5 & Google ADK\nBuilt for the AI Hackathon 2026"

    # --- Slide 2: Project Overview & Tech Stack ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Project Overview & Tech Stack"

    content = slide.placeholders[1]
    content.text = "An automated feedback analysis system that turns raw reviews into actionable intelligence."
    
    p = content.text_frame.add_paragraph()
    p.text = "Key Components:"
    p.level = 1
    
    components = [
        "LLM: Gemini 2.5 Flash (Hyper-fast, large context)",
        "Framework: Google Agent Development Kit (ADK)",
        "Deployment: Google Cloud Run (Serverless Scaling)",
        "Interface: Modern Chatbot UI with Sentiment Visualization"
    ]
    
    for comp in components:
        p = content.text_frame.add_paragraph()
        p.text = comp
        p.level = 2

    # --- Slide 3: Outcomes & Impact ---
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Outcomes & Future Impact"

    content = slide.placeholders[1]
    outcomes = [
        "100% Automated Feedback Triage: Instant sentiment detection.",
        "Concise Summarization: Distills long reviews into 1-paragraph insights.",
        "Response Optimization: Drafts professional replies in the brand voice.",
        "Scalability: Ready to handle thousands of reviews via Cloud Run."
    ]
    
    for outcome in outcomes:
        p = content.text_frame.add_paragraph()
        p.text = outcome
        p.level = 1

    # Save the presentation
    output_path = "Customer_Feedback_Intelligence_AI.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
